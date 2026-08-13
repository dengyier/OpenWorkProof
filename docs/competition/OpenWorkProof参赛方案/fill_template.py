"""Fill the official GOAI template with OpenWorkProof content."""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt

SRC = "/Users/molin/Downloads/2e567d1a-99c1-45ce-8a0f-3d36d11f3314 (1).pptx"
DST = "/Users/molin/Project/openWorkProof-scope-bound-verification-v03/docs/competition/OpenWorkProof参赛方案/OpenWorkProof初赛方案-官方模板.pptx"
IMG = "/Users/molin/Project/openWorkProof-scope-bound-verification-v03/docs/competition/OpenWorkProof参赛方案/images"

shutil.copy(SRC, DST)
prs = Presentation(DST)


def set_text(shape, text, size=None, bold=None, color=None):
    """Replace shape text keeping the first run's formatting baseline."""
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        if para.runs:
            run = para.runs[0]
            run.text = text if first else ""
            first = False
            if size:
                run.font.size = Pt(size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color
            # clear extra runs
            for extra in para.runs[1:]:
                extra.text = ""
        else:
            para.text = text if first else ""
            first = False


def find_shape(slide, contains):
    for sh in slide.shapes:
        if sh.has_text_frame and contains in sh.text_frame.text:
            return sh
    return None


def add_bullets(slide, left, top, width, height, items, size=14):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = lvl
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.name = "Microsoft YaHei"
    return tb


def replace_pictures(slide, image_paths):
    """Replace PICTURE shapes in order with given images."""
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    for pic, path in zip(pics, image_paths):
        left, top, w, h = pic.left, pic.top, pic.width, pic.height
        pic._element.getparent().remove(pic._element)
        slide.shapes.add_picture(path, left, top, width=w, height=h)


def delete_shape(slide, contains):
    for sh in list(slide.shapes):
        if sh.has_text_frame and contains in sh.text_frame.text:
            sh._element.getparent().remove(sh._element)
            return True
    return False


# ---------------------------------------------------------------- Slide 1
s = prs.slides[0]
set_text(find_shape(s, "初赛方案 PPT"), "初赛方案 PPT：Agent 工作契约与可验证执行协议", 36, True)
set_text(find_shape(s, "Agent Infra"), "Agent Infra · 新智基座  |  OpenWorkProof", 18, False)
delete_shape(s, "Datawhale")

# ---------------------------------------------------------------- Slide 2  P0
s = prs.slides[1]
p0 = {
    "【在此填写项目名称】": "OpenWorkProof —— Agent 工作契约与可验证执行协议",
    "【描述真实场景与核心痛点】": "企业 Agent 交付「完成」不可验证、不可审计、不可离线复核：合规审计一件工作 3–5 天，只能信厂商自陈",
    "【概述端到端解决方案】": "AgentTeams 多 Agent 协同 + 8 个 OWP Skill + 契约/授权/证据/验收协议层，让每步执行可验签、可离线复核",
    "【列 1–2 个关键差异化优势】": "① Proof-Carrying Work：交付物自带证据链 ② 机器验证闭环：owp-verifier 确定性复核（非 LLM 自证）③ 同源生态：修复 AgentScope 自身 bug",
    "【说明复用与迁移价值】": "Apache-2.0 · 8 个可复用 Skill · 25 个 MCP 工具 · 协议可迁移到任意多 Agent 系统（SaaS/金融/研发）",
    "【说明当前完成度与里程碑】": "AgentTeams 真实接入跑通（分派证据存档）· 3056 测试 0 失败 · AgentScope #2239 可复现 demo · 初赛进行中",
}
for key, val in p0.items():
    sh = find_shape(s, key)
    if sh:
        set_text(sh, val, 13, False)

# ---------------------------------------------------------------- Slide 5
s = prs.slides[4]
replace_pictures(s, [f"{IMG}/scene_scenarios.png", f"{IMG}/scene_diff.png"])

# ---------------------------------------------------------------- Slide 7
s = prs.slides[6]
replace_pictures(s, [f"{IMG}/architecture.png"])

# ---------------------------------------------------------------- Slide 9
s = prs.slides[8]
replace_pictures(s, [f"{IMG}/multiagent_team.png", f"{IMG}/multiagent_roles.png"])

# ---------------------------------------------------------------- Slide 11
s = prs.slides[10]
replace_pictures(s, [f"{IMG}/skills_s1s8.png"])

# ---------------------------------------------------------------- Slide 13
s = prs.slides[12]
replace_pictures(s, [f"{IMG}/eng_verification.png", f"{IMG}/eng_evidence.png",
                     f"{IMG}/eng_observability.png", f"{IMG}/eng_security.png"])

# ---------------------------------------------------------------- Slide 15
s = prs.slides[14]
delete_shape(s, "建议覆盖")
add_bullets(s, 0.7, 2.2, 11.8, 4.6, [
    ("可复用成果：协议层（WorkOrder/Grant/Receipt/Acceptance 四原语）、双 AgentTeams 适配器（agt 管理面 + Matrix 执行面）、供应链镜像工具（prepare_context / convert_docker_archive）、8 个 Skill 与 25 个 MCP 工具", 0),
    ("接口契约与文档：README 中英双语、MCP_SERVER.md、offline-verification.md、docs/status.md（真值边界）、docs/competition/ 全套参赛材料", 0),
    ("开源协议：Apache-2.0；代码仓库 github.com/dengyier/OpenWorkProof（2,283+ 测试，交付验证 M1–M4 已签署）", 0),
    ("第三方依赖：Python 标准库为主 + cryptography / pydantic / FastMCP；AgentTeams（Apache-2.0）仅做集成，不拥有其信任根", 0),
    ("生态：对阿里云官方 Skills 对齐设计；AgentTeams 为必选协同基点，已实证接入", 0),
], size=15)

# ---------------------------------------------------------------- Slide 17
s = prs.slides[16]
replace_pictures(s, [f"{IMG}/plan_milestones.png", f"{IMG}/plan_risks.png"])

# ---------------------------------------------------------------- Slide 19
s = prs.slides[18]
delete_shape(s, "可从以下方面介绍团队基本情况")
add_bullets(s, 0.8, 1.6, 11.5, 5.6, [
    ("核心团队", 0),
    ("· 创始人 / 技术 Owner —— OpenWorkProof 发起人，负责协议架构与总体设计；20 年从业经验，擅长跨领域工程与产品化", 1),
    ("· 董浩宇 博士 —— 验证方法学；主导交付验证计划（M1–M4）、Ed25519 签署与验收权威分离设计", 1),
    ("· 邓海波 —— 供应链与镜像交付；负责 candidate inventory、wheelhouse/deb 闭包与 Docker 供应链门", 1),
    ("· 龙胜海 —— 协议与安全边界；负责授权链、攻击矩阵（C0+A1–A18）与权限模型评审", 1),
    ("团队分工", 0),
    ("· 协议设计 / 执行证据 / 审计账本：创始人 + 龙胜海", 1),
    ("· 交付验证 / 测试门 / 签署：董浩宇", 1),
    ("· 供应链 / 容器 / 发布工程：邓海波", 1),
    ("· AgentTeams 接入 / 复赛 Demo / 商业试点：全团队 + 顾问支持", 1),
    ("团队成果", 0),
    ("· OpenWorkProof（Apache-2.0）：v0.1–v0.4 全协议栈，3056 测试 0 失败，独立双审 READY，发布候选已推送远端 main", 1),
    ("· AgentTeams（hiclaw）真实接入：双适配器 14 测试通过，Matrix 程序化闭环验证，Manager 自动分派真实修复任务（证据存档）", 1),
], size=14)

prs.save(DST)
print("saved:", DST)

# verify
prs2 = Presentation(DST)
print("slides:", len(prs2.slides))
